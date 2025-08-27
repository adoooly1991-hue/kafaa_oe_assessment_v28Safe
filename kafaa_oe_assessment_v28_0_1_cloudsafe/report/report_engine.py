
import json, argparse, os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.dml.color import RGBColor

from .charts import product_chart, waterfall, ecrs_funnel

ASSETS = Path('assets')
DEFAULT_TEMPLATE = ASSETS/'Kafaa_Template.pptx'
DEFAULT_BRAND = ASSETS/'kafaa_logo.png'

def add_textbox(slide, x, y, w, h, text, size=18, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame; tf.clear()
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = text; run.font.size = Pt(size); run.font.bold = bold
    p.alignment = align
    return box

def add_kpi_card(slide, x, y, w, h, label, value, color=(27,159,151)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(I(x)._EMUS), int(I(y)._EMUS), int(I(w)._EMUS), int(I(h)._EMUS))
    fill = shape.fill; fill.solid(); fill.fore_color.rgb = RGBColor(242,250,249)
    line = shape.line; line.color.rgb = RGBColor(*color)
    tf = shape.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = label; r.font.size = Pt(12); r.font.bold = True
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = value; r2.font.size = Pt(18)
    return shape

def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists(): return None
    if w and h: return slide.shapes.add_picture(path, I(x), I(y), width=I(w), height=I(h))
    return slide.shapes.add_picture(path, I(x), I(y))

def cover_slide(prs, payload, brand_mode):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_textbox(slide, 0.8, 1.2, 8.0, 1.0, "Operational Excellence Assessment", 36, True)
    add_textbox(slide, 0.8, 2.1, 8.0, 0.8, payload['client'].get('name_en', 'Client'), 20)
    # logos
    brand_logo = str(DEFAULT_BRAND)
    client_logo = payload['client'].get('logo')
    add_image(slide, brand_logo, 0.8, 0.3, 2.2, 0.9)
    if brand_mode == 'co_brand' and client_logo:
        add_image(slide, client_logo, 8.0-2.2, 0.3, 2.0, 0.8)
    if brand_mode == 'white_label' and client_logo:
        add_image(slide, client_logo, 0.8, 0.3, 2.2, 0.9)

def agenda_slide(prs, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph(); p.text = b; p.level = 0

def methodology_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "Assessment Methodology", 24, True)
    steps = ["Data Collection","Financial Study","Product Selection","VSM (Current/Future/Ideal)","Observations (PQCDSM)","Impact & Charter","Exports"]
    x = 0.9
    for s in steps:
        box = add_kpi_card(slide, x, 1.3, 2.1, 1.1, s, "", color=(88,89,91))
        x += 2.2

def financials_slide(prs, payload):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "Financial Targets", 24, True)
    f = payload.get('financials',{})
    add_kpi_card(slide, 0.8, 1.3, 3.0, 1.2, "Target Profit (SAR)", f"{f.get('target_profit_sar','—'):,}")
    add_kpi_card(slide, 3.9, 1.3, 3.0, 1.2, "Cost Reduction Target (SAR)", f"{f.get('cost_reduction_target_sar','—'):,}")
    add_kpi_card(slide, 0.8, 2.7, 3.0, 1.2, "Revenue (SAR)", f"{f.get('revenue_sar','—'):,}")

def selection_slide(prs, payload):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "Champion Product Selection", 24, True)
    sel = payload.get('product_selection',{})
    add_kpi_card(slide, 0.8, 1.2, 3.2, 1.2, "Champion product", sel.get('champion','—'))
    # charts (if provided)
    charts = []
    series = sel.get('series') or []
    if series:
        p1 = product_chart(series, "Sales Value", "SAR")
        p2 = product_chart(series, "Sales Volume", "Units")
        charts.extend([p1,p2])
    x=0.8
    for p in charts[:2]:
        add_image(slide, p, x, 2.6, 3.3, 2.0); x+=3.4

def vsm_summary_slide(prs, payload):
    # notes for vsm
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "VSM Summary — Current / Future / Ideal", 24, True)
    cur, fut, ide = payload['vsm']['current'], payload['vsm']['future'], payload['vsm'].get('ideal',{})
    add_kpi_card(slide, 0.8, 1.3, 3.0, 1.3, "Current", f"TP={cur.get('touch_points','—')}, LT={cur.get('lead_time_days','—')} days, VA%={cur.get('va_pct','—')}")
    add_kpi_card(slide, 3.9, 1.3, 3.0, 1.3, "Future", f"TP={fut.get('touch_points','—')}, LT={fut.get('lead_time_days','—')} days, VA%={fut.get('va_pct','—')}")
    if ide:
        add_kpi_card(slide, 0.8, 2.8, 3.0, 1.3, "Ideal", f"TP={ide.get('touch_points','—')}, LT={ide.get('lead_time_hrs','—')} hrs, VA%={ide.get('va_pct','—')}")

    try:
        notes = f"VSM: Current LT {cur.get('lead_time_days','—')} d, Future LT {fut.get('lead_time_days','—')} d."
        _set_notes(slide, notes)
    except Exception:
        pass

def waterfalls_trio_slide(prs, payload):
    # notes for waterfalls
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "Impact — Waterfalls", 24, True)
    wf = payload.get('muda',{}).get('quantification',{}).get('waterfalls',{})
    imgs = []
    if wf.get('cost'):
        imgs.append(waterfall(wf['cost'], "Cost Savings (SAR/yr)"))
    if wf.get('cash'):
        imgs.append(waterfall(wf['cash'], "Frozen Cash (SAR)"))
    if wf.get('lost_opportunity'):
        imgs.append(waterfall(wf['lost_opportunity'], "Sales Opportunity (SAR/yr)"))
    x=0.8
    for p in imgs:
        add_image(slide, p, x, 1.3, 2.7, 2.0); x+=2.9

    try:
        costs = wf.get('cost',[]); cash = wf.get('cash',[]); lost = wf.get('lost_opportunity',[])
        note = f"Savings items: {len(costs)} • Frozen cash items: {len(cash)} • Sales opp items: {len(lost)}."
        _set_notes(slide, note)
    except Exception:
        pass

def ecrs_funnel_slide(prs, payload):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "ECRS Funnel", 24, True)
    m = payload.get('muda',{})
    n_from = int(m.get('raw_count', 110))
    n_to = int(m.get('ecrs_final', 13))
    img = ecrs_funnel(n_from, n_to)
    add_image(slide, img, 0.8, 1.3, 3.0, 2.2)

def photos_grid_slide(prs, payload):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "Gemba Photos", 24, True)
    imgs = (payload.get('images') or {}).get('gemba') or []
    x, y = 0.8, 1.2; col = 0
    for p in imgs[:6]:
        add_image(slide, p, x, y, 2.5, 1.8)
        x += 2.7; col += 1
        if col % 3 == 0: x = 0.8; y += 2.0

def assumptions_slide(prs, payload):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Assumptions"
    body = slide.shapes.placeholders[1].text_frame; body.clear()
    for a in payload.get('assumptions', []):
        p = body.add_paragraph(); p.text = f"• {a}"; p.level = 0

def disclaimer_slide(prs, payload):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Disclaimer"
    body = slide.shapes.placeholders[1].text_frame; body.clear()
    p = body.add_paragraph(); p.text = payload.get('disclaimer_en', 'This report is based on client-provided data.'); p.level = 0

def signoff_slide(prs, payload, brand_mode):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.6, 8.0, 0.6, "Sign-off & Acceptance", 24, True)
    # signature lines
    add_textbox(slide, 0.8, 2.2, 3.2, 0.3, "Kafaa Representative", 12)
    add_textbox(slide, 4.2, 2.2, 3.2, 0.3, f"{payload['client'].get('name_en','Client')} Representative", 12)
    # logos
    add_image(slide, str(DEFAULT_BRAND), 0.8, 0.3, 2.2, 0.9)
    if brand_mode in ('co_brand','white_label') and payload['client'].get('logo'):
        add_image(slide, payload['client']['logo'], 8.0-2.2, 0.3, 2.0, 0.8)

def build_pptx(payload_path:str, blueprint_path:str, out_path:str, brand_mode:str='kafaa'):
    payload = json.loads(Path(payload_path).read_text(encoding='utf-8'))
    bp = json.loads(Path(blueprint_path).read_text(encoding='utf-8').replace("\t"," ")) if blueprint_path.endswith('.json') else None
    if blueprint_path.endswith('.yaml') or blueprint_path.endswith('.yml'):
        import yaml
        bp = yaml.safe_load(Path(blueprint_path).read_text(encoding='utf-8'))
    template = bp.get('template') if bp else str(DEFAULT_TEMPLATE)
    prs = Presentation(template)
    # slides
    for s in (bp.get('slides') if bp else []):
        t = s.get('type')
        if t=='cover': cover_slide(prs, payload, brand_mode)
        elif t=='executive_summary': executive_summary_slide(prs, payload)
        elif t=='agenda': agenda_slide(prs, s.get('bullets',[]))
        elif t=='methodology': methodology_slide(prs)
        elif t=='financials': financials_slide(prs, payload)
        elif t=='product_selection': selection_slide(prs, payload)
        elif t=='vsm_summary': vsm_summary_slide(prs, payload)
        elif t=='waterfalls_trio': waterfalls_trio_slide(prs, payload)
        elif t=='ecrs_funnel': ecrs_funnel_slide(prs, payload)
        elif t=='photos_grid': photos_grid_slide(prs, payload)
        elif t=='assumptions': assumptions_slide(prs, payload)
        elif t=='disclaimer': disclaimer_slide(prs, payload)
        elif t=='signoff': signoff_slide(prs, payload, brand_mode)
        elif t=='appendix_obs': appendix_obs_slide(prs, payload)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path

if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('--payload', required=True, help='Path to report_payload.json')
    ap.add_argument('--blueprint', default='report/blueprint.yaml', help='Path to blueprint YAML/JSON')
    ap.add_argument('--out', default='exports/Assessment_Report.pptx')
    ap.add_argument('--brand', default='kafaa', choices=['kafaa','co_brand','white_label'])
    args = ap.parse_args()
    path = build_pptx(args.payload, args.blueprint, args.out, args.brand)
    print('Wrote:', path)


from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.lib.utils import ImageReader

def build_pdf(payload_path:str, blueprint_path:str, out_path:str, brand_mode:str='kafaa'):
    payload = json.loads(Path(payload_path).read_text(encoding='utf-8'))
    W, H = landscape(A4)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(out_path, pagesize=landscape(A4))
    # Cover
    c.setFont("Helvetica-Bold", 22); c.drawString(2*cm, H-3*cm, "Operational Excellence Assessment")
    c.setFont("Helvetica", 14); c.drawString(2*cm, H-4*cm, payload.get('client',{}).get('name_en','Client'))
    try:
        brand_logo = ImageReader(str(DEFAULT_BRAND))
        c.drawImage(brand_logo, 2*cm, H-2.5*cm, width=4*cm, height=1.5*cm, mask='auto')
    except Exception: pass
    if brand_mode in ('co_brand','white_label'):
        cl = payload.get('client',{}).get('logo')
        if cl and Path(cl).exists():
            try:
                c.drawImage(ImageReader(cl), W-6*cm, H-2.5*cm, width=4*cm, height=1.5*cm, mask='auto')
            except Exception: pass
    c.showPage()
    # Financial targets
    f = payload.get('financials',{})
    c.setFont("Helvetica-Bold", 18); c.drawString(2*cm, H-2.5*cm, "Financial Targets")
    c.setFont("Helvetica", 12)
    c.drawString(2*cm, H-3.8*cm, f"Target Profit (SAR): {f.get('target_profit_sar','—'):,}")
    c.drawString(2*cm, H-4.8*cm, f"Cost Reduction Target (SAR): {f.get('cost_reduction_target_sar','—'):,}")
    c.drawString(2*cm, H-5.8*cm, f"Revenue (SAR): {f.get('revenue_sar','—'):,}")
    c.showPage()
    # Waterfalls
    wf = payload.get('muda',{}).get('quantification',{}).get('waterfalls',{})
    c.setFont("Helvetica-Bold", 18); c.drawString(2*cm, H-2.5*cm, "Impact — Waterfalls")
    x = 2*cm; y = H-11*cm
    def _draw_img(p):
        nonlocal x
        if p and Path(p).exists():
            c.drawImage(ImageReader(p), x, y, width=8*cm, height=6*cm, preserveAspectRatio=True, mask='auto')
            x += 9*cm
    # generate images via charts.py helpers
    try:
        imgs = []
        if wf.get('cost'): imgs.append(waterfall(wf['cost'], "Cost Savings (SAR/yr)"))
        if wf.get('cash'): imgs.append(waterfall(wf['cash'], "Frozen Cash (SAR)"))
        if wf.get('lost_opportunity'): imgs.append(waterfall(wf['lost_opportunity'], "Sales Opportunity (SAR/yr)"))
        for p in imgs: _draw_img(p)
    except Exception: pass
    c.showPage()
    # Sign-off
    c.setFont("Helvetica-Bold", 18); c.drawString(2*cm, H-2.5*cm, "Sign-off & Acceptance")
    c.setFont("Helvetica", 12); c.drawString(2*cm, H-4*cm, "Kafaa Representative: ______________________   Date: __________")
    c.drawString(2*cm, H-5*cm, f"{payload.get('client',{}).get('name_en','Client')} Representative: ______________________   Date: __________")
    c.save()
    return out_path


def executive_summary_slide(prs, payload):
    # Executive Summary narrative
    narrative = payload.get('exec_narrative','')
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "Executive Summary", 24, True)
    # Key KPIs
    f = payload.get('financials',{})
    vsm = payload.get('vsm',{})
    cur, fut = vsm.get('current',{}), vsm.get('future',{})
    add_kpi_card(slide, 0.8, 1.4, 2.8, 1.2, "Cost Reduction Target (SAR)", f"{f.get('cost_reduction_target_sar','—'):,}")
    add_kpi_card(slide, 3.0, 1.4, 2.8, 1.2, "Frozen Cash (SAR)", f"{int(sum([i.get('value',0) for i in (payload.get('muda',{}).get('quantification',{}).get('waterfalls',{}).get('cash',[]))])):,}")
    add_kpi_card(slide, 5.2, 1.4, 2.8, 1.2, "Sales Opportunity (SAR/yr)", f"{int(sum([i.get('value',0) for i in (payload.get('muda',{}).get('quantification',{}).get('waterfalls',{}).get('lost_opportunity',[]))])):,}")
    # Narrative box
    if narrative:
        add_textbox(slide, 0.8, 3.9, 7.2, 0.9, narrative[:250], 12, False)
    # VSM delta cards
    try:
        lt_delta = float(cur.get('lead_time_days',0)) - float(fut.get('lead_time_days',0))
    except Exception:
        lt_delta = 0
    add_kpi_card(slide, 0.8, 2.8, 2.8, 1.2, "Lead Time (Current)", f"{cur.get('lead_time_days','—')} days")
    add_kpi_card(slide, 3.0, 2.8, 2.8, 1.2, "Lead Time (Future)", f"{fut.get('lead_time_days','—')} days")
    add_kpi_card(slide, 5.2, 2.8, 2.8, 1.2, "Reduction", f"{lt_delta:.0f} days")
    # Champion product
    _set_notes(slide, (narrative or 'Executive summary.') )
    # Champion product
    add_kpi_card(slide, 0.8, 4.2, 7.2, 1.2, "Champion Product", payload.get('product_selection',{}).get('champion','—'))


def appendix_obs_slide(prs, payload):
    obs = payload.get('observations') or []
    rag_srcs = payload.get('rag_sources') or []
    if not obs and not rag_srcs:
        return
    # First slide for observations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Appendix — Observations"
    body = slide.shapes.placeholders[1].text_frame; body.clear()
    if obs:
        for i, o in enumerate(obs[:10]):
            p = body.add_paragraph(); p.text = f"• {o[:220]}"; p.level = 0
    # If we have more, add another slide
    if len(obs) > 10:
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Appendix — Observations (cont.)"
        b2 = slide2.shapes.placeholders[1].text_frame; b2.clear()
        for o in obs[10:20]:
            p = b2.add_paragraph(); p.text = f"• {o[:220]}"; p.level = 0
    # Sources slide
    if rag_srcs:
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "Appendix — Source Snippets"
        b3 = slide3.shapes.placeholders[1].text_frame; b3.clear()
        for i, s in enumerate(rag_srcs[:8]):
            p = b3.add_paragraph(); p.text = f"• {s[:220]}"; p.level = 0


def _set_notes(slide, text):
    try:
        ns = slide.notes_slide
    except Exception:
        ns = slide.notes_slide if hasattr(slide, "notes_slide") else slide.part.notes_slide
    try:
        tf = slide.notes_slide.notes_text_frame
        tf.clear(); tf.text = text
    except Exception:
        try:
            slide.notes_slide.notes_text_frame.text = text
        except Exception:
            pass

# tag: Executive Summary narrative

# notes: exec

from PIL import Image, ImageDraw, ImageFont

def build_preview(payload_path:str, out_path:str, brand_mode:str='kafaa'):
    payload = json.loads(Path(payload_path).read_text(encoding='utf-8'))
    W,H = 1200, 675  # 16:9 preview
    im = Image.new('RGB',(W,H),(245,247,250))
    d = ImageDraw.Draw(im)
    # Title
    title = "Operational Excellence Assessment"
    client = payload.get('client',{}).get('name_en','Client')
    subt = f"Client: {client}"
    try:
        f1 = ImageFont.truetype('DejaVuSans-Bold.ttf', 44)
        f2 = ImageFont.truetype('DejaVuSans.ttf', 28)
    except:
        f1 = f2 = None
    d.text((60,60), title, fill=(30,30,30), font=f1)
    d.text((60,120), subt, fill=(60,60,60), font=f2)
    # KPIs
    f = payload.get('financials',{})
    kpis = [
        ("Cost Reduction Target (SAR)", f"{f.get('cost_reduction_target_sar','—'):,}"),
        ("Frozen Cash (SAR)", f"{int(sum([i.get('value',0) for i in (payload.get('muda',{}).get('quantification',{}).get('waterfalls',{}).get('cash',[]))])):,}"),
        ("Sales Opportunity (SAR/yr)", f"{int(sum([i.get('value',0) for i in (payload.get('muda',{}).get('quantification',{}).get('waterfalls',{}).get('lost_opportunity',[]))])):,}")
    ]
    y = 200
    for label, val in kpis:
        d.rectangle([50,y, 1150,y+80], outline=(197,31,45), width=2)
        d.text((70,y+18), f"{label}: {val}", fill=(40,40,40), font=f2)
        y += 100
    # logos
    try:
        brand_logo = str(DEFAULT_BRAND)
        if Path(brand_logo).exists():
            lb = Image.open(brand_logo).convert("RGBA").resize((280,110))
            im.paste(lb, (60,H-150), lb)
    except Exception: pass
    cl = payload.get('client',{}).get('logo')
    if brand_mode in ('co_brand','white_label') and cl and Path(cl).exists():
        try:
            lc = Image.open(cl).convert("RGBA").resize((260,100))
            im.paste(lc, (W-320,H-150), lc)
        except Exception: pass
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    im.save(out_path, format="PNG")
    return out_path

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib.utils import ImageReader

def build_reviewer_checklist_pdf(payload_path:str, approvals:dict, out_path:str, sig_kafaa:str=None, sig_client:str=None):
    payload = json.loads(Path(payload_path).read_text(encoding='utf-8'))
    c = canvas.Canvas(out_path, pagesize=A4)
    W, H = A4
    # Header
    c.setFont("Helvetica-Bold", 16)
    c.drawString(2*cm, H-2.5*cm, "Reviewer Checklist — Operational Excellence Assessment")
    c.setFont("Helvetica", 11)
    client = payload.get('client',{}).get('name_en','Client')
    c.drawString(2*cm, H-3.2*cm, f"Client: {client}")
    # Observations approvals
    obs = payload.get('observations', []) or []
    y = H-4.0*cm
    c.setFont("Helvetica-Bold", 12); c.drawString(2*cm, y, "Observations — Decisions")
    y -= 0.6*cm
    c.setFont("Helvetica", 10)
    for i, o in enumerate(obs[:20]):
        c.drawString(2*cm, y, f"{i+1}. {o[:90]}")
        c.drawString(14*cm, y, approvals.get(str(i), "Pending"))
        y -= 0.6*cm
        if y < 3*cm:
            c.showPage(); y = H-3*cm
    # Signatures
    c.setFont("Helvetica-Bold", 12)
    y = 3.0*cm
    c.drawString(2*cm, y+1.2*cm, "Kafaa Representative")
    c.rect(2*cm, y, 6*cm, 1*cm)
    if sig_kafaa and Path(sig_kafaa).exists():
        try:
            c.drawImage(ImageReader(sig_kafaa), 2.1*cm, y, width=5.8*cm, height=1.0*cm, preserveAspectRatio=True, mask='auto')
        except Exception: pass
    c.drawString(11*cm, y+1.2*cm, f"{client} Representative")
    c.rect(11*cm, y, 6*cm, 1*cm)
    if sig_client and Path(sig_client).exists():
        try:
            c.drawImage(ImageReader(sig_client), 11.1*cm, y, width=5.8*cm, height=1.0*cm, preserveAspectRatio=True, mask='auto')
        except Exception: pass
    c.save()
    return out_path

from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def attach_reviewer_checklist_to_pptx(pptx_in:str, payload_path:str, approvals:dict, out_path:str, sig_kafaa:str=None, sig_client:str=None):
    payload = json.loads(Path(payload_path).read_text(encoding='utf-8'))
    prs = Presentation(pptx_in)
    # Title slide for appendix
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts)>5 else prs.slide_layouts[1])
    try:
        slide.shapes.title.text = "Appendix — Reviewer Checklist"
    except Exception:
        box = slide.shapes.add_textbox(I(0.8), I(0.5), I(8.0), I(0.8))
        box.text_frame.text = "Appendix — Reviewer Checklist"
    # Observations decisions
    obs = payload.get('observations', []) or []
    # Create subsequent slides if needed
    idx = 0
    while idx < len(obs):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        try:
            s.shapes.title.text = "Observations — Decisions"
            body = s.shapes.placeholders[1].text_frame; body.clear()
        except Exception:
            body = s.shapes.add_textbox(I(0.8), I(1.0), I(8.0), I(4.5)).text_frame; body.clear()
        for i in range(idx, min(idx+10, len(obs))):
            stxt = f"{i+1}. {obs[i][:200]} — {approvals.get(str(i),'Pending')}"
            p = body.add_paragraph(); p.text = stxt; p.level = 0
        idx += 10
    # Signature slide
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    def add_text(slide, x,y,w,h, text, size=14, bold=False):
        tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h)).text_frame
        tb.clear(); p = tb.paragraphs[0]; r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    add_text(s2, 0.8, 0.6, 8.0, 0.5, "Reviewer Signatures", 20, True)
    add_text(s2, 0.8, 1.2, 3.0, 0.3, "Kafaa Representative", 12)
    add_text(s2, 4.5, 1.2, 3.0, 0.3, f"{payload.get('client',{}).get('name_en','Client')} Representative", 12)
    # signature boxes
    try:
        if sig_kafaa and Path(sig_kafaa).exists():
            s2.shapes.add_picture(sig_kafaa, I(0.8), I(1.5), width=I(3.0), height=I(0.9))
        else:
            add_text(s2, 0.8, 1.5, 3.0, 0.9, " ", 12)
        if sig_client and Path(sig_client).exists():
            s2.shapes.add_picture(sig_client, I(4.5), I(1.5), width=I(3.0), height=I(0.9))
        else:
            add_text(s2, 4.5, 1.5, 3.0, 0.9, " ", 12)
    except Exception:
        pass
    prs.save(out_path)
    return out_path
