from components.boot import boot
mode = boot()

import streamlit as st
import json, os
from pathlib import Path
from report.report_engine import build_pptx, build_pdf, build_preview, build_reviewer_checklist_pdf, attach_reviewer_checklist_to_pptx
from components.coach_agent import recommend_with_trace
, io, os, plotly.graph_objects as go, tempfile, pandas as pd
from pptx import Presentation
from pptx.util import Inches
from components.export_plus import load_template, title_slide, section_slide, text_block, add_impact_slide, table_slide, thankyou_slide, kpi_card, vision_grid_slide
from components.roi_bridge import rollup_impact, series_to_waterfall_data
st.set_page_config(page_title="Export", layout="wide")
st.title("📤 Export PPTX (Kafaa template + KPIs + Vision grid)")
profile = st.session_state.get('profile_key','default')
template_path = "assets/Kafaa_Template.pptx"; logo_path = "assets/kafaa_logo.png"; client_logo = st.session_state.get('client_logo_path')
if not client_logo and 'serb' in (st.session_state.get('client_name','').lower()):
    if os.path.exists('assets/serb_logo.png'):
        client_logo = 'assets/serb_logo.png'
client_name = (st.session_state.get('client_name','') or '').lower()
if 'serb' in client_name and os.path.exists('assets/Serb_Template.pptx'):
    template_path = 'assets/Serb_Template.pptx'
elif profile.startswith('defense') and os.path.exists('assets/Kafaa_Template_Defense.pptx'):
    template_path = 'assets/Kafaa_Template_Defense.pptx'
def pick_kpis(ss):
    champion = ss.get("champion_product"); champ_info = ""
    try:
        df = ss.get("product_table")
        if champion is not None and isinstance(df, pd.DataFrame) and "Product" in df.columns:
            row = df[df["Product"]==champion].head(1)
            if not row.empty:
                margin = row.iloc[0].get("MarginPerUnit") or row.iloc[0].get("Gross-Margin / Unit  (SAR)")
                sales = row.iloc[0].get("TotalSales") or row.iloc[0].get("Sales (SAR)")
                champ_info = f"Champion: {champion} — Margin/unit: {margin}, Sales: {sales}"
    except Exception:
        pass
    takt = ss.get("takt_sec")
    if not takt:
        shift = ss.get("shift_sec") or 8*3600; demand = ss.get("demand_per_shift") or None
        if demand: takt = shift/float(demand)
    ct = None
    if ss.get("ct_table"):
        try: ct = max([r.get("ct",0) for r in ss.get("ct_table")])
        except Exception: pass
    if not ct:
        prod = ss.get("coach_mode_results",{}).get("production",{}).get("answers",{}); ct = prod.get("ct_bn")
    prod = ss.get("coach_mode_results",{}).get("production",{}).get("answers",{}); fpy = prod.get("fpy_line")
    inv = ss.get("financials_inventory") or 0; rate = ss.get("carrying_rate_pct") or 25; carrying = inv * (float(rate)/100.0)
    return {"takt": takt, "ct_bn": ct, "fpy": fpy, "carrying": carrying, "champ_info": champ_info}
def build_impact_images(state):
    data = rollup_impact(state); imgs = {}; tmpdir = tempfile.mkdtemp()
    def draw(items, title, key):
        if not items: return
        x,y,m = series_to_waterfall_data(items)
        fig = go.Figure(go.Waterfall(x=x, y=y, measure=m))
        fig.update_layout(title=title, height=360, width=600, margin=dict(l=8,r=8,t=40,b=8))
        path = os.path.join(tmpdir, f"{key}.png"); fig.write_image(path); imgs[key] = path
    draw(data.get("savings_items"), "Savings (SAR/yr)", "wf_savings")
    draw(data.get("frozen_cash_items"), "Frozen Cash (SAR)", "wf_frozen")
    draw(data.get("sales_oppty_items"), "Sales Opportunity (SAR/yr)", "wf_sales")
    return imgs
def harvest_vision(ss):
    out = []; store = ss.get("vision_evidence") or {}
    for stage, items in store.items():
        for it in items:
            if os.path.exists(it.get("path","")): out.append((f"{stage}: {it.get('label','')}", it["path"]))
    return out
if st.button("Generate PPTX"):
    prs = load_template(template_path)
    title = f"Kafaa OE — Assessment ({st.session_state.get('factory_name','Factory')})"; subtitle = f"Year: {st.session_state.get('year','')} — Client: {st.session_state.get('client_name','')} — Project: {st.session_state.get('project_name','')}"
    cover = title_slide(prs, title, subtitle=subtitle, logo_path=logo_path)
try:
    from pptx.util import Inches as _I
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches as _I

def _defense_header(slide):
    try:
        w = slide.part.presentation.slide_width; h = slide.part.presentation.slide_height
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, int(h*0.05))
        bar.fill.solid();
        # use tokenized accent
        from pptx.dml.color import RGBColor as _RGB
        hexcol = accent_color_hex().lstrip('#'); r=int(hexcol[0:2],16); g=int(hexcol[2:4],16); b=int(hexcol[4:6],16)
        bar.fill.fore_color.rgb = _RGB(r,g,b); bar.line.fill.background()
        # icon strip
        import os
        for i,icon in enumerate(["assets/icon_uav.png","assets/icon_shield.png","assets/icon_chip.png"]):
            if os.path.exists(icon):
                slide.shapes.add_picture(icon, _I(0.4 + i*0.9), _I(0.08), width=_I(0.6))
    except Exception:
        pass

    if client_logo: cover.shapes.add_picture(client_logo, _I(0.5), _I(0.2), width=_I(1.2))
    # accent rectangle on cover if defense
    if profile.startswith('defense'):
        w = cover.part.presentation.slide_width; h = cover.part.presentation.slide_height
        shp = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, int(h*0.05))
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(11,61,145)  # defense navy
        shp.line.fill.background()
except Exception:
    pass  # accent is optional
    k = pick_kpis(st.session_state)
    if any([v for kk,v in k.items() if kk!='champ_info']):
        kpi_card(cover, 0.6, 4.0, 2.3, 1.2, "Takt (sec)", f"{k.get('takt'):.1f}" if k.get('takt') else "—")
        kpi_card(cover, 3.0, 4.0, 2.3, 1.2, "Bottleneck CT (sec)", f"{k.get('ct_bn'):.1f}" if k.get('ct_bn') else "—")
        kpi_card(cover, 5.4, 4.0, 2.3, 1.2, "FPY (%)", f"{k.get('fpy'):.1f}" if k.get('fpy') else "—")
        kpi_card(cover, 7.8, 4.0, 2.3, 1.2, "Carrying Cost (SAR/yr)", f"{k.get('carrying'):,.0f}")
        if k.get('champ_info'):
            cover.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(9.5), Inches(0.6)).text_frame.text = k.get('champ_info')
    # Charter slide\n    slide = section_slide(prs, "VSM Charter", logo_path=logo_path)\n    chart_lines = []\n    crt = st.session_state.get('cost_reduction_target_sar'); invr = st.session_state.get('inv_reduction_target_sar')\n    quick = st.session_state.get('quick_ratio'); curr = st.session_state.get('current_ratio')\n    champ = st.session_state.get('champion_product')\n    chart_lines.append(f"Champion Product: {champ or '—'}")\n    chart_lines.append(f"Cost Reduction Target: {crt:,.0f} SAR")\n    chart_lines.append(f"Inventory Reduction Target: {invr:,.0f} SAR")\n    chart_lines.append(f"Quick Ratio: {quick if quick is not None else '—'} | Current Ratio: {curr if curr is not None else '—'}")\n    scope = st.session_state.get('charter_scope',''); obj = st.session_state.get('charter_objectives','')\n    if scope: chart_lines.append('Scope: ' + scope)\n    if obj: chart_lines.append('Objectives: ' + obj)\n    text_block(slide, chart_lines)\n\n    obs = st.session_state.get("observations", [])
    if obs:
        slide = section_slide(prs, "Observations (PQCDSM)", logo_path=logo_path)
        paras = []
        for o in obs:
            paras.append(f"{o.get('section','')} — {o.get('title','')}"); paras.append(o.get("text",""))
        text_block(slide, paras)
    cm = st.session_state.get("coach_mode_results", {})
    if cm:
        slide = section_slide(prs, "Coach Mode — Top Wastes by Stage", logo_path=logo_path)
        lines = []
        for kstage,v in cm.items():
            wastes = ", ".join([f"{w} (s{sc})" for w,sc in v.get("top_wastes",[])])
            if wastes: lines.append(f"{v.get('title','')} → {wastes}")
        if lines: text_block(slide, lines)
    df = st.session_state.get("countermeasures_df")
    if isinstance(df, pd.DataFrame) and not df.empty:
        table_slide(prs, "Countermeasures Plan", df, logo_path=logo_path)
    imgs = build_impact_images(st.session_state)
    if imgs: add_impact_slide(prs, imgs, logo_path=logo_path)
    vision = harvest_vision(st.session_state)
    if vision: vision_grid_slide(prs, "Gemba Photo Evidence", vision, logo_path=logo_path)
    thankyou_slide(prs, "Thank you", logo_path=logo_path)
    import io; bio = io.BytesIO(); prs.save(bio)
    fname = f"Kafaa_OE_Assessment_{st.session_state.get('client_name','Client')}_{st.session_state.get('project_name','Project')}.pptx".replace(' ','_')
log("Export PPTX", {"file":"pptx"})
    st.download_button("Download PPTX", bio.getvalue(), file_name=fname)
else:
    st.info("Click **Generate PPTX** to export with KPI cards and photo evidence grid.")


def add_before_after_slide(prs):
    from pptx.util import Inches as I, Pt
    from pptx.dml.color import RGBColor
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title if slide.shapes.title else slide.shapes.add_textbox(I(0.6), I(0.3), I(9), I(1)).text_frame
    try:
        title.text = "Before vs After (Targets)"
    except:
        tf = slide.shapes.add_textbox(I(0.6), I(0.3), I(9), I(1)).text_frame; tf.text = "Before vs After (Targets)"
    # gather numbers
    takt = st.session_state.get("takt_sec"); ct = None
    prod = st.session_state.get("coach_mode_results",{}).get("production",{}).get("answers",{})
    ct = prod.get("ct_bn") or st.session_state.get("ct_bottleneck_sec")
    fpy = prod.get("fpy_line") or st.session_state.get("fpy_line")
    bench_fpy = 98.0
    try:
        from components.bench import load_profile
        prof = load_profile(st.session_state.get("profile_key","default"))
        bench_fpy = float(prof.get("kpis",{}).get("fpy_best_practice_pct", 98))
    except Exception:
        pass
    inv = float(st.session_state.get("financials_inventory") or 0.0)
    target_turns = float(st.session_state.get("target_turns") or 8.0)
    # crude current turns proxy if COGS available: COGS/avg inv * 365; else '--'
    try:
        current_turns = (float(st.session_state.get("financials_cogs") or 0.0) / (inv if inv else 1.0)) * 365.0/365.0
    except Exception:
        current_turns = None
    target_savings = float(st.session_state.get("cost_reduction_target_sar") or 0.0)
    pace = st.session_state.get("pace_sums", {}) or {}
    now_next = float(pace.get("Now",0.0) + pace.get("Next",0.0))

    rows = [
        ("Bottleneck CT vs Takt (s)", f"{ct or '--'} → {takt or '--'}"),
        (f"FPY vs Best Practice ({bench_fpy}%)", f"{fpy or '--'}% → {bench_fpy}%"),
        ("Inventory Turns", f"{current_turns:.1f} → {target_turns}") if current_turns else ("Inventory Turns","— → {target_turns}"),
        ("Savings vs Target (SAR/yr)", f"{now_next:,.0f} → {target_savings:,.0f}")
    ]
    x0, y0 = I(0.8), I(1.6)
    for i,(k,v) in enumerate(rows):
        box = slide.shapes.add_textbox(x0, y0 + I(0.8*i), I(11), I(0.7)).text_frame
        box.text = f"{k}: {v}"


def add_audit_slide(prs):
    from pptx.util import Inches as I
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(I(0.6), I(0.3), I(9), I(1)).text_frame
    tf.text = "Audit Trail"
    log = st.session_state.get("audit_log", [])
    y = 1.2
    for item in log[-12:]:
        t = slide.shapes.add_textbox(I(0.8), I(y), I(11), I(0.5)).text_frame
        t.text = f"{item.get('ts')} — {item.get('event')}"
        y += 0.6


def add_signoff_slide(prs):
    from pptx.util import Inches as I
    from pptx.dml.color import RGBColor
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(I(0.6), I(0.3), I(9), I(0.8)).text_frame
    tf.text = "Sign-off and Acceptance"
    # signature lines
    y = 3.5
    slide.shapes.add_textbox(I(0.8), I(2.2), I(5), I(0.4)).text_frame.text = "For Kafaa:"
    slide.shapes.add_textbox(I(6.0), I(2.2), I(5), I(0.4)).text_frame.text = "For Client:"
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(72*0.8), int(72*3.0), int(72*4.5), int(2))
    line1.fill.solid(); line1.fill.fore_color.rgb = RGBColor(120,120,120); line1.line.fill.background()
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(72*6.0), int(72*3.0), int(72*4.5), int(2))
    line2.fill.solid(); line2.fill.fore_color.rgb = RGBColor(120,120,120); line2.line.fill.background()
    slide.shapes.add_textbox(I(0.8), I(3.1), I(5), I(0.3)).text_frame.text = "Name / Title / Signature / Date"
    slide.shapes.add_textbox(I(6.0), I(3.1), I(5), I(0.3)).text_frame.text = "Name / Title / Signature / Date"

from components.heatmap_util import generate_and_store_heatmap
_generate_heat = True
try:
    generate_and_store_heatmap()
except Exception:
    pass


def _memo_excerpts(max_items=2):
    memos = st.session_state.get("gemba_memos") or []
    out = []
    for m in memos[:max_items]:
        if len(m) > 240: m = m[:237] + "..."
        out.append(m)
    return out

def _photo_paths(max_items=2):
    p = st.session_state.get("photo_evidence") or []
    return [rec.get("path") for rec in p[:max_items] if isinstance(rec, dict) and rec.get("path")]

def add_observation_slides(prs):
    from pptx.util import Inches as I, Pt
    obs_list = st.session_state.get("observations")
    if not obs_list:
        # fallback: build a single slide from RAG answer if present
        ans = st.session_state.get("rag_last_answer") or ""
        if not ans: return
        obs_list = [ans]
    low = st.session_state.get("obs_low_confidence", False)
    for i, obs in enumerate(obs_list[:6]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
        slide.shapes.title.text = f"Observation {i+1}" + (" — LOW CONFIDENCE" if low else "")
        body = slide.shapes.placeholders[1].text_frame
        body.clear(); p = body.paragraphs[0]; p.text = obs[:2000]
        # confidence bar
        try:
            from components.obs_conf import bulk_confidences
            confs = bulk_confidences([obs])
            c = confs[0] if confs else {"measured":0,"inferred":0}
            total = max(1, c['measured']+c['inferred'])
            m_ratio = c['measured']/total
            from pptx.util import Inches as I
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(72*0.8), int(72*4.4), int(72*5.0), int(72*0.15))
            fill = bar.fill; fill.solid(); fill.fore_color.rgb = RGBColor(208,212,217)
            mbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(72*0.8), int(72*4.4), int(72*5.0*m_ratio), int(72*0.15))
            mbar.fill.solid(); mbar.fill.fore_color.rgb = RGBColor(27,159,151)
            cap = slide.shapes.add_textbox(I(0.8), I(4.55), I(5.0), I(0.2)).text_frame; cap.text = f"Confidence: Measured {c['measured']} / Inferred {c['inferred']}"
        except Exception:
            pass

        # memo excerpts
        ex = _memo_excerpts()
        if ex:
            box = slide.shapes.add_textbox(I(0.8), I(4.8), I(5.4), I(1.3)).text_frame
            box.text = "Memos:"
            for e in ex:
                r = box.add_paragraph(); r.text = "• " + e
        # thumbnails
        imgs = _photo_paths()
        x = 6.6; y = 4.6
        for j, path in enumerate(imgs):
            try:
                slide.shapes.add_picture(path, I(x), I(y), width=I(2.1))
                x += 2.2
            except Exception:
                pass


def _safe(n, default=None):
    import streamlit as st
    return st.session_state.get(n, default)

def _waterfalls_from_state():
    pace = _safe('pace_sums', {}) or {}
    invA = float(_safe('abc_A_value', 0) or 0)
    invB = float(_safe('abc_B_value', 0) or 0)
    invC = float(_safe('abc_C_value', 0) or 0)
    sales_ct = float(_safe('lost_sales_ct_gap', 0) or 0)
    sales_chg = float(_safe('lost_sales_changeover', 0) or 0)
    sales_q = float(_safe('lost_sales_quality', 0) or 0)
    return {
        "cost": [{"label":"Now","value":float(pace.get('Now',0))},{"label":"Next","value":float(pace.get('Next',0))},{"label":"Later","value":float(pace.get('Later',0))},{"label":"Total","value":0,"measure":"total"}],
        "cash": [{"label":"A","value":invA},{"label":"B","value":invB},{"label":"C","value":invC},{"label":"Total","value":0,"measure":"total"}],
        "lost_opportunity": [{"label":"CT gaps","value":sales_ct},{"label":"Changeover","value":sales_chg},{"label":"Quality","value":sales_q},{"label":"Total","value":0,"measure":"total"}],
    }

def _payload_from_state(include_rag=False):
    import streamlit as st
    client = _safe('client_profile', {"name_en": _safe('client_name','Client'), "logo": _safe('client_logo_path', None)})
    payload = {
        "brand_mode": _safe('brand_mode','kafaa'),
        "client": client,
        "financials": {
            "target_profit_sar": float(_safe('target_profit_sar', 0) or 0),
            "cost_reduction_target_sar": float(_safe('cost_reduction_target_sar', 0) or 0),
            "revenue_sar": float(_safe('revenue_sar', 0) or 0),
        },
        "product_selection": {
            "champion": _safe('champion_product','—'),
            "series": _safe('product_series', [])
        },
        "vsm": {
            "current": {"touch_points": int(_safe('tp_current', _safe('touch_points_current', 0)) or 0),
                        "lead_time_days": float(_safe('lt_days_current', 0) or 0),
                        "va_pct": float(_safe('va_pct_current', 0) or 0)},
            "future": {"touch_points": int(_safe('tp_future', _safe('touch_points_future', 0)) or 0),
                       "lead_time_days": float(_safe('lt_days_future', 0) or 0),
                       "va_pct": float(_safe('va_pct_future', 0) or 0)},
            "ideal": {"touch_points": int(_safe('tp_ideal', 0) or 0),
                      "lead_time_hrs": float(_safe('lt_hrs_ideal', 0) or 0),
                      "va_pct": float(_safe('va_pct_ideal', 0) or 0)},
        },
        "muda": {
            "raw_count": int(_safe('muda_raw_count', 0) or 0),
            "ecrs_final": int(_safe('muda_ecrs_final', 0) or 0),
            "quantification": {"waterfalls": _waterfalls_from_state()}
        },
        "images": {"gemba": [rec.get('path') for rec in (_safe('photo_evidence') or []) if isinstance(rec, dict) and rec.get('path')]},
        "assumptions": _safe('assumptions_list', ["8h shift, 5d/week","SAR currency","Lead time excludes weekends"]),
        "disclaimer_en": _safe('disclaimer_en', 'This report is based on client-provided data and observations during Gemba walks.')
    }
    
# Optional RAG appendix
if include_rag:
    rag_ans = _safe('rag_last_answer', '')
    obs_list = _safe('observations', []) or ([rag_ans] if rag_ans else [])
    payload['observations'] = [o for o in obs_list if o]
    payload['rag_sources'] = _safe('rag_last_sources', [])

payload['exec_narrative'] = _exec_narrative()
payload['coach_trace'] = [r.get('because','') for r in recommend_with_trace(st.session_state.get('vsm_steps') or [])]
return payload

st.subheader("Full Report (PPTX) — Kafaa Template")
brand = st.selectbox("Brand mode", ["kafaa","co_brand","white_label"], index=["kafaa","co_brand","white_label"].index(_safe('brand_mode','kafaa')))
also_pdf = st.checkbox("Also create PDF (lite)")
append_checklist = st.checkbox("Append Reviewer Checklist to PPTX (as appendix)")
disabled = (int(score*100) < int(min_ok))
if disabled:
    st.warning("Completeness below threshold. Please fill the missing items or lower the threshold.")
if st.button("Generate full report", type="primary") and not disabled:
    payload = _payload_from_state(include_rag=include_rag)
    Path("exports").mkdir(exist_ok=True, parents=True)
    Path("report").mkdir(exist_ok=True, parents=True)
    live_payload = Path("report/live_payload.json")
    live_payload.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    out_pptx = f"exports/Assessment_Report_{payload['client'].get('name_en','Client').replace(' ','_')}.pptx"
    path = build_pptx(str(live_payload), "report/blueprint.yaml", out_pptx, brand)
    st.success(f"Created {path}")
    st.download_button("Download PPTX", data=open(out_pptx,'rb').read(), file_name=Path(out_pptx).name)

    # Build preview PNG + manifest entry
    preview_path = out_pptx.replace('.pptx', '.preview.png')
    build_preview(str(live_payload), preview_path, brand)
    manifest_p = Path('exports/manifest.json')
    man = {}
    if manifest_p.exists():
        try: man = json.loads(manifest_p.read_text(encoding='utf-8'))
        except Exception: man = {}
    wsname = st.session_state.get('active_workspace_name') or st.session_state.get('client_name','workspace')
    rec = {"pptx": out_pptx, "preview": preview_path, "ts": __import__('datetime').datetime.now().isoformat(timespec='seconds')}
    hist = man.get(wsname, [])
    hist.append(rec)
    man[wsname] = hist[-5:]  # keep last 5
    manifest_p.write_text(json.dumps(man, ensure_ascii=False, indent=2), encoding='utf-8')


    if append_checklist:
        out_pptx2 = out_pptx.replace('.pptx','_with_appendix.pptx')
        approvals = st.session_state.get("approvals", {})
        attach_reviewer_checklist_to_pptx(out_pptx, str(live_payload), approvals, out_pptx2, st.session_state.get('sig_kafaa_path'), st.session_state.get('sig_client_path'))
        st.success(f"Appended checklist to PPTX → {out_pptx2}")
        st.download_button("Download PPTX (with Appendix)", data=open(out_pptx2,'rb').read(), file_name=Path(out_pptx2).name)

    if also_pdf:
        out_pdf = out_pptx.replace('.pptx','.pdf')
        p2 = build_pdf(str(live_payload), "report/blueprint.yaml", out_pdf, brand)
        st.success(f"Created {p2}")
        st.download_button("Download PDF", data=open(out_pdf,'rb').read(), file_name=Path(out_pdf).name)


def _is_set(x):
    return x not in (None, "", [], {}, 0, 0.0)

def _completeness():
    req = {
        "Financials: target_profit_sar": _safe('target_profit_sar', None),
        "Financials: cost_reduction_target_sar": _safe('cost_reduction_target_sar', None),
        "Financials: revenue_sar": _safe('revenue_sar', None),
        "Selection: champion_product": _safe('champion_product', None),
        "VSM: TP current & future": all([_safe('tp_current', None), _safe('tp_future', None)]),
        "VSM: LT current & future": all([_safe('lt_days_current', None), _safe('lt_days_future', None)]),
        "Waterfalls: cost": (_waterfalls_from_state().get('cost') or []),
        "Waterfalls: cash": (_waterfalls_from_state().get('cash') or []),
        "Waterfalls: lost_opportunity": (_waterfalls_from_state().get('lost_opportunity') or []),
    }
    ok = sum(1 for v in req.values() if _is_set(v))
    total = len(req)
    score = ok/total if total else 1.0
    missing = [k for k,v in req.items() if not _is_set(v)]
    return score, missing

st.markdown("### Completeness")
min_ok = st.slider("Minimum completeness required", 0, 100, 75, help="Block export if below this percent")
score, missing = _completeness()
st.progress(score)
st.caption(f"Data completeness: {int(score*100)}% (min {min_ok}%)")
if missing:
    with st.expander("Missing items (click to expand)"):
        for m in missing:
            st.write("• ", m)

st.markdown("---")
st.subheader("Appendix")
include_rag = st.checkbox("Include RAG‑grounded appendix (latest draft + sources)")


def _exec_narrative():
    steps = st.session_state.get("vsm_steps") or []
    recs = recommend_with_trace(steps)
    parts = []
    # Pull key metrics
    takt = st.session_state.get("takt_sec", None)
    bn = None
    try:
        # find bottleneck from steps
        bn = max(steps, key=lambda s: float(s.get("ct_sec",0))) if steps else None
    except Exception:
        bn = None
    if takt and bn:
        parts.append(f"Bottleneck CT {bn.get('ct_sec')}s vs takt {takt}s at {bn.get('name')}.")
    chg = st.session_state.get("changeover_min", None)
    if chg: parts.append(f"Changeover ~{chg} min.")
    fpy = st.session_state.get("fpy_current", None)
    if fpy: parts.append(f"FPY {fpy}%.")
    # Recommendations summary
    if recs:
        top = "; ".join([r['action'] for r in recs[:3]])
        parts.append(f"Top actions: {top}.")
    return " ".join(parts)[:250]


st.markdown("---")
st.subheader("Reviewer Checklist")
sig1 = st.text_input("Kafaa signature image path (optional)", value=st.session_state.get("sig_kafaa_path",""))
sig2 = st.text_input("Client signature image path (optional)", value=st.session_state.get("sig_client_path",""))
if st.button("Generate Reviewer Checklist (PDF)"):
    payload = _payload_from_state(include_rag=include_rag)
    Path("exports").mkdir(exist_ok=True, parents=True)
    live_payload = Path("report/live_payload.json"); live_payload.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    approvals = st.session_state.get("approvals", {})
    out_pdf = f"exports/Reviewer_Checklist_{payload['client'].get('name_en','Client').replace(' ','_')}.pdf"
    build_reviewer_checklist_pdf(str(live_payload), approvals, out_pdf, sig1 or None, sig2 or None)
    st.success(f"Created {out_pdf}")
    st.download_button("Download Checklist PDF", data=open(out_pdf,'rb').read(), file_name=Path(out_pdf).name)


st.markdown("---")
st.subheader("Live preview")
if st.button("Refresh live preview"):
    payload = _payload_from_state(include_rag=include_rag)
    Path("report").mkdir(exist_ok=True, parents=True)
    live_payload = Path("report/live_payload.json")
    live_payload.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    # Use current Brand selection
    preview_out = f"exports/Live_Preview_{payload['client'].get('name_en','Client').replace(' ','_')}.png"
    build_preview(str(live_payload), preview_out, brand)
    st.image(preview_out, caption="Live preview", width=720)
