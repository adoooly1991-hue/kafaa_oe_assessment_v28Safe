
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
def load_template(path):
    try: return Presentation(path)
    except Exception: return Presentation()
def _brand_logo(slide, logo_path, width_in=1.1):
    try: slide.shapes.add_picture(logo_path, Inches(9.0), Inches(0.2), width=Inches(width_in))
    except Exception: pass
def title_slide(prs, title, subtitle="", logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = title; slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    if len(slide.placeholders) > 1: slide.placeholders[1].text = subtitle or datetime.now().strftime("%Y-%m-%d")
    if logo_path: _brand_logo(slide, logo_path); return slide
def kpi_card(slide, x_in, y_in, w_in, h_in, label, value, sub=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    shp.fill.solid(); shp.line.color.rgb = shp.fill.fore_color.rgb
    tf = shp.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = str(value); p.font.size = Pt(28); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(12); p2.alignment = PP_ALIGN.CENTER
    if sub: p3 = tf.add_paragraph(); p3.text = sub; p3.font.size = Pt(10); p3.alignment = PP_ALIGN.CENTER
    return shp
def section_slide(prs, title, logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    if slide.shapes.title: slide.shapes.title.text = title; slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    if logo_path: _brand_logo(slide, logo_path); return slide
def text_block(slide, paragraphs, left=0.6, top=1.5, width=10.0, height=4.7):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.clear()
    for i, p in enumerate(paragraphs):
        par = tf.add_paragraph() if i>0 else tf.paragraphs[0]
        par.text = p; par.font.size = Pt(14); par.level = 0
    return tb
def add_impact_slide(prs, images_dict, logo_path=None):
    slide = section_slide(prs, "Impact — Savings / Frozen Cash / Sales Opportunity", logo_path=logo_path)
    x = Inches(0.5); y = Inches(1.5); w = Inches(3.0); h = Inches(3.0)
    keys = ["wf_savings","wf_frozen","wf_sales"]
    for i,k in enumerate(keys):
        img = images_dict.get(k)
        if img: slide.shapes.add_picture(img, x + Inches(3.2*i), y, width=w, height=h)
    return slide
def table_slide(prs, title, df, logo_path=None, max_rows=18):
    slide = section_slide(prs, title, logo_path=logo_path)
    rows = min(len(df.index)+1, max_rows); cols = len(df.columns)
    if rows<2 or cols<1: return slide
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.5), Inches(5.0)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0,j); cell.text = str(col); cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.size = Pt(12)
    for i in range(rows-1):
        for j, col in enumerate(df.columns):
            val = str(df.iloc[i, j]) if i < len(df.index) else ""
            cell = table.cell(i+1, j); cell.text = val; cell.text_frame.paragraphs[0].font.size = Pt(11)
    return slide
def vision_grid_slide(prs, title, images_labeled, logo_path=None, cols=3):
    slide = section_slide(prs, title, logo_path=logo_path)
    if not images_labeled: return slide
    x0, y0, w, h = 0.5, 1.5, 3.0, 2.2
    for idx, (label, path) in enumerate(images_labeled[:9]):
        r = idx // cols; c = idx % cols
        slide.shapes.add_picture(path, Inches(x0 + c*(w+0.2)), Inches(y0 + r*(h+0.4)), width=Inches(w), height=Inches(h))
        tb = slide.shapes.add_textbox(Inches(x0 + c*(w+0.2)), Inches(y0 + r*(h+0.4) + h + 0.05), Inches(w), Inches(0.3))
        tf = tb.text_frame; tf.text = label; tf.paragraphs[0].font.size = Pt(10)
    return slide
def thankyou_slide(prs, text="Thank you", logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    if slide.shapes.title: slide.shapes.title.text = text
    if logo_path: _brand_logo(slide, logo_path)
    return slide
